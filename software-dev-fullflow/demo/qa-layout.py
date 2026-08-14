# 程序化布局 QA：检查 PPT 元素是否超出边界/重叠/文本可能溢出
# 无法渲染图片，用坐标做启发式检查
from pptx import Presentation
from pptx.util import Emu

path = r"c:\Users\34239\Documents\GitHub\mutil_agent\software-dev-fullflow\demo\GOAI-赛道三-方案.pptx"
prs = Presentation(path)
SW, SH = prs.slide_width / 914400, prs.slide_height / 914400  # inches
print(f"Slide size: {SW:.2f}x{SH:.2f}\n")

def inches(v): return v / 914400 if v is not None else None

issues = []
for si, slide in enumerate(prs.slides, 1):
    shapes = []
    for sh in slide.shapes:
        if sh.shape_type is None or not sh.has_text_frame:
            continue
        x, y, w, h = inches(sh.left), inches(sh.top), inches(sh.width), inches(sh.height)
        if x is None: continue
        shapes.append((sh, x, y, w, h))
    # 边界检查
    for sh, x, y, w, h in shapes:
        # 跳过全宽背景条/圆点装饰（允许贴边）
        name = sh.shape_type
        if x < -0.02 or y < -0.02:
            issues.append(f"S{si} '{sh.text_frame.text[:20]}' 左/上越界 x={x:.2f} y={y:.2f}")
        if x + w > SW + 0.02:
            issues.append(f"S{si} '{sh.text_frame.text[:20]}' 右越界 x+w={x+w:.2f} > {SW}")
        if y + h > SH + 0.02:
            issues.append(f"S{si} '{sh.text_frame.text[:20]}' 下越界 y+h={y+h:.2f} > {SH}")
    # 卡片重叠检查（文本块两两，忽略装饰矩形）
    txt = [s for s in shapes if s[0].has_text_frame and s[0].text_frame.text.strip()]
    for i in range(len(txt)):
        for j in range(i+1, len(txt)):
            a, b = txt[i], txt[j]
            # 同一卡片内标题+正文不算重叠；跳过完全相同的 x,y 起始
            if abs(a[1]-b[1])<0.02 and abs(a[2]-b[2])<0.02: continue
            ax,ay,aw,ah = a[1],a[2],a[3],a[4]
            bx,by,bw,bh = b[1],b[2],b[3],b[4]
            ox = max(0, min(ax+aw,bx+bw)-max(ax,bx))
            oy = max(0, min(ay+ah,by+bh)-max(ay,by))
            if ox>0.05 and oy>0.05:
                issues.append(f"S{si} 文本重叠: '{a[0].text_frame.text[:12]}' × '{b[0].text_frame.text[:12]}' ({ox:.2f}x{oy:.2f})")

if issues:
    print("发现潜在问题:")
    for i in issues: print("  -", i)
else:
    print("✅ 无边界/重叠问题")

# 文本溢出启发式：估算文本量 vs 框大小（粗略，中文字符宽 ~ fontSize pt）
print("\n--- 潜在文本过密检查（仅宽高比极端者）---")
for si, slide in enumerate(prs.slides, 1):
    for sh in slide.shapes:
        if not sh.has_text_frame: continue
        t = sh.text_frame.text.strip()
        if not t: continue
        w, h = inches(sh.width), inches(sh.height)
        if w is None or h is None: continue
        fs = None
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size: fs = r.font.size.pt; break
            if fs: break
        if not fs: fs = 11
        # 每行可容纳中文字符 ≈ w*96/fs*1.0 (字符宽~pt)，行数 = h*96/(fs*1.6)
        per_line = max(int(w*96/fs), 1)
        est_lines = max(1, int(h*96/(fs*1.6)))
        need_lines = 0
        for line in t.split("\n"):
            need_lines += max(1, (len(line)+per_line-1)//per_line)
        if need_lines > est_lines*1.5:
            print(f"S{si} '{t[:25]}' 需~{need_lines}行 > 框~{est_lines}行 (fontSize {fs}, w{w:.2f} h{h:.2f})")
print("\nQA 完成")
