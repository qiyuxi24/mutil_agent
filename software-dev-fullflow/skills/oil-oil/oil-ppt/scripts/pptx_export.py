#!/usr/bin/env python3
"""Hybrid PPTX export from canonical final HTML and its rendered DOM."""
from __future__ import annotations

import base64
import hashlib
import json
import os
import re
import shutil
import subprocess
import tempfile
import time
import urllib.request
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Callable

from cdp_validate import WebSocket, _stop_browser, _wait_for_devtools_port
from theme import PALETTES, TYPOGRAPHY, validate_theme

SLIDE_WIDTH_PX = 1920
SLIDE_HEIGHT_PX = 1080
SLIDE_WIDTH_EMU = 12_192_000
SLIDE_HEIGHT_EMU = 6_858_000
COVERAGE_SCHEMA = "oil-ppt.pptx-editability/v2"


def require_python_pptx() -> Any:
    """Import python-pptx only for export."""
    try:
        import pptx
    except (ImportError, OSError) as error:
        raise SystemExit("PPTX export requires python-pptx: python3 -m pip install python-pptx") from error
    return pptx


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _cdp_command(socket: WebSocket, command_id: int, method: str, params: dict | None = None) -> dict:
    socket.send_json({"id": command_id, "method": method, **({"params": params} if params else {})})
    while True:
        message = socket.recv_json()
        if message.get("id") != command_id:
            continue
        if "error" in message:
            raise RuntimeError(f"Chrome DevTools {method} failed: {message['error']}")
        return message.get("result") or {}


def _browser_session(chrome: str, html_path: Path) -> tuple[subprocess.Popen, tempfile.TemporaryDirectory, WebSocket]:
    profile_ctx = tempfile.TemporaryDirectory(prefix="oil-ppt-pptx-cdp-")
    profile = Path(profile_ctx.name)
    command = [
        chrome, "--headless", "--no-sandbox", "--allow-file-access-from-files",
        "--disable-background-networking", "--disable-component-update", "--disable-default-apps",
        "--disable-features=PaintHolding,RenderDocument", "--disable-sync", "--force-color-profile=srgb",
        "--force-device-scale-factor=1", "--hide-scrollbars", "--metrics-recording-only", "--no-first-run",
        "--window-size=1920,1080", f"--user-data-dir={profile}", "--remote-debugging-port=0", html_path.resolve().as_uri(),
    ]
    process = subprocess.Popen(command, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, start_new_session=True)
    try:
        deadline, port_file = time.monotonic() + 15, profile / "DevToolsActivePort"
        port = _wait_for_devtools_port(port_file, deadline, "Chrome did not expose a DevTools port for PPTX export.")
        pages: list[dict] = []
        while time.monotonic() < deadline and not pages:
            try:
                with urllib.request.urlopen(f"http://127.0.0.1:{port}/json/list", timeout=2) as response:
                    pages = [item for item in json.load(response) if item.get("type") == "page"]
            except OSError:
                time.sleep(.05)
        if not pages:
            raise RuntimeError("Chrome opened no inspectable page for PPTX export.")
        target = html_path.resolve().as_uri()
        page = next((item for item in pages if item.get("url", "").startswith(target)), pages[0])
        socket = WebSocket(page["webSocketDebuggerUrl"])
        socket.sock.settimeout(20)
        return process, profile_ctx, socket
    except Exception:
        _stop_browser(process, profile)
        profile_ctx.cleanup()
        raise


def _layout_expression(slide_index: int) -> str:
    """Extract visible leaf text/images and suppress only safe native candidates."""
    return f"""(async () => {{
      await document.fonts?.ready;
      await Promise.all([...document.images].map(image => image.complete ? Promise.resolve() : new Promise(done => {{ image.onload = image.onerror = done; }})));
      document.documentElement.style.cssText += ';width:1920px!important;height:1080px!important';
      document.body.style.cssText += ';width:1920px!important;height:1080px!important;overflow:hidden!important';
      const viewport = document.querySelector('.deck-viewport,.slide-preview-viewport');
      const shell = document.querySelector('.deck-stage-shell,.slide-preview-shell');
      const stage = document.querySelector('.deck-stage,.slide-preview-stage');
      if (!stage) throw new Error('canonical HTML has no deck stage');
      const overview = document.querySelector('[data-deck-overview]');
      if (overview && !overview.hidden) document.querySelector('[data-deck-overview-close]')?.click();
      if (overview && !overview.hidden) throw new Error('overview must be closed before PPTX capture');
      const exportChrome = [...document.querySelectorAll('.deck-counter,.progress-bar,.next-preview,.deck-overview-toggle,.deck-overview')];
      exportChrome.forEach(node => node.style.display = 'none');
      if (viewport) viewport.style.cssText += ';position:fixed!important;inset:0!important;width:1920px!important;height:1080px!important';
      if (shell) shell.style.cssText += ';position:relative!important;width:1920px!important;height:1080px!important';
      stage.style.cssText += ';position:absolute!important;left:0!important;top:0!important;width:1920px!important;height:1080px!important;transform:none!important';
      const slides = [...document.querySelectorAll('.oil-slide')], slide = slides[{slide_index}];
      if (!slide) throw new Error('missing slide index {slide_index}');
      if (stage.querySelectorAll(':scope > .oil-slide').length !== slides.length) throw new Error('all slides must be restored to stage before PPTX capture');
      slides.forEach((node, index) => {{ node.classList.toggle('active', index === {slide_index}); node.style.cssText += index === {slide_index} ? ';display:block!important;visibility:visible!important;opacity:1!important;transition:none!important' : ';display:none!important;visibility:hidden!important;opacity:0!important;transition:none!important'; }});
      await new Promise(done => requestAnimationFrame(() => requestAnimationFrame(done)));
      const visible = node => {{ if (!node?.getClientRects().length) return false; const s=getComputedStyle(node); return s.display !== 'none' && s.visibility !== 'hidden' && Number(s.opacity || 1) > .001; }};
      const boxFor = node => {{ const b=node.getBoundingClientRect(), root=slide.getBoundingClientRect(); return {{x:b.left-root.left,y:b.top-root.top,width:b.width,height:b.height}}; }};
      const styledAncestor = node => {{
        for (let parent=node.parentElement; parent && parent !== slide; parent=parent.parentElement) {{
          const s=getComputedStyle(parent), overflow=[s.overflow,s.overflowX,s.overflowY];
          const clipped=overflow.some(value => ['hidden','clip'].includes(value))
            || !['none',''].includes(s.clipPath) || !['none',''].includes(s.maskImage)
            || parseFloat(s.borderTopLeftRadius || '0') > .01 || parseFloat(s.borderTopRightRadius || '0') > .01
            || parseFloat(s.borderBottomLeftRadius || '0') > .01 || parseFloat(s.borderBottomRightRadius || '0') > .01;
          if (clipped || s.transform !== 'none' || s.filter !== 'none' || Number(s.opacity || 1) < .999) return true;
        }}
        return false;
      }};
      const textElements = [...slide.querySelectorAll('*')].filter(node => visible(node) && !['script','style','svg','path'].includes(node.tagName.toLowerCase()))
        .filter(node => (node.textContent || '').replace(/\\s+/g,' ').trim())
        .filter(node => ![...node.querySelectorAll('*')].some(child => visible(child) && (child.textContent || '').replace(/\\s+/g,' ').trim()));
      const text = textElements.map((node, index) => {{
        const s=getComputedStyle(node), value=(node.textContent || '').replace(/\\s+/g,' ').trim(), box=boxFor(node);
        const generated = pseudo => {{ const content=getComputedStyle(node,pseudo).content; return content && !['none','normal','""',"''"].includes(content); }};
        const zero = raw => Math.abs(parseFloat(raw || '0')) < .01;
        const native=box.width > .5 && box.height > .5
          && s.transform === 'none' && s.filter === 'none' && s.opacity === '1'
          && s.textShadow === 'none' && s.textDecorationLine === 'none'
          && s.textTransform === 'none' && s.writingMode === 'horizontal-tb'
          && s.display !== 'list-item' && !styledAncestor(node)
          && ['normal','0px'].includes(s.letterSpacing) && zero(s.textIndent)
          && [s.paddingTop,s.paddingRight,s.paddingBottom,s.paddingLeft].every(zero)
          && [s.borderTopWidth,s.borderRightWidth,s.borderBottomWidth,s.borderLeftWidth].every(zero)
          && (!s.backgroundColor || s.backgroundColor === 'rgba(0, 0, 0, 0)')
          && !generated('::before') && !generated('::after');
        const item={{index,selector:`${{node.tagName.toLowerCase()}}:nth-leaf(${{index+1}})`,text:value,tag:node.tagName.toLowerCase(),box,fontFamily:s.fontFamily,fontSize:s.fontSize,fontWeight:s.fontWeight,fontStyle:s.fontStyle,color:s.color,textAlign:s.textAlign,lineHeight:s.lineHeight,letterSpacing:s.letterSpacing,whiteSpace:s.whiteSpace,native,reason:native ? '' : 'computed text style is not safely representable as native PowerPoint text'}};
        if (native) node.style.visibility='hidden'; return item;
      }}).filter(item => item.box.width > .5 && item.box.height > .5);
      const images=[...slide.querySelectorAll('img')].filter(visible).map((node,index) => {{
        const s=getComputedStyle(node), box=boxFor(node), src=node.currentSrc || node.src, portable=/^data:image\\/(png|jpeg|gif|webp);base64,/i.test(src) || /^file:.*\\.(png|jpe?g|gif|webp)(?:[?#]|$)/i.test(src), styled=box.width > .5 && box.height > .5 && s.transform === 'none' && s.filter === 'none' && s.opacity === '1' && s.mixBlendMode === 'normal' && s.boxShadow === 'none' && s.clipPath === 'none' && s.maskImage === 'none' && (s.borderRadius === '0px' || s.borderRadius === '0') && !styledAncestor(node) && ['fill','contain','cover'].includes(s.objectFit || 'fill'), native=portable && styled;
        const item={{index,selector:`img:nth-of-type(${{index+1}})`,src,box,objectFit:s.objectFit || 'fill',objectPosition:s.objectPosition || '50% 50%',naturalWidth:node.naturalWidth,naturalHeight:node.naturalHeight,native,reason:native ? '' : (!portable ? 'rendered image source cannot be embedded reliably by python-pptx' : 'computed image style is not safely representable as native PowerPoint media')}};
        if (native) node.style.visibility='hidden'; return item;
      }});
      await new Promise(done => requestAnimationFrame(() => requestAnimationFrame(done)));
      return {{slideId:slide.dataset.slideId || '',title:slide.dataset.title || '',text,images,slideBox:boxFor(slide),chromeHidden:exportChrome.every(node => getComputedStyle(node).display === 'none')}};
    }})()"""


def collect_render_layers(chrome: str, html_path: Path, deck: dict, background_dir: Path) -> list[dict]:
    """Capture one background plus final-DOM layers for each manifest slide."""
    paths = deck.get("slides") if isinstance(deck, dict) else None
    if not isinstance(paths, list) or not paths:
        raise RuntimeError("deck manifest has no slides for PPTX export")
    process, profile_ctx, socket = _browser_session(chrome, html_path)
    profile, results, command_id = Path(profile_ctx.name), [], 1
    try:
        _cdp_command(socket, command_id, "Page.enable"); command_id += 1
        _cdp_command(socket, command_id, "Runtime.enable"); command_id += 1
        for index, _ in enumerate(paths):
            evaluated = _cdp_command(socket, command_id, "Runtime.evaluate", {"expression": _layout_expression(index), "awaitPromise": True, "returnByValue": True}); command_id += 1
            remote = evaluated.get("result") or {}
            if remote.get("subtype") == "error" or "exceptionDetails" in evaluated:
                raise RuntimeError(f"Could not inspect rendered slide {index + 1}: {remote.get('description')}")
            layout = remote.get("value")
            if not isinstance(layout, dict) or not layout.get("slideId"):
                raise RuntimeError(f"Rendered slide {index + 1} is missing its data-slide-id.")
            screenshot = _cdp_command(socket, command_id, "Page.captureScreenshot", {"format": "png", "fromSurface": True, "clip": {"x": 0, "y": 0, "width": SLIDE_WIDTH_PX, "height": SLIDE_HEIGHT_PX, "scale": 1}}); command_id += 1
            destination = background_dir / f"slide-{index + 1:03d}.png"
            destination.write_bytes(base64.b64decode(screenshot["data"]))
            layout["background"] = str(destination)
            results.append(layout)
    finally:
        socket.close(); _stop_browser(process, profile); profile_ctx.cleanup()
    return results


def classify_rendered_slide(layout: dict) -> dict:
    """Pure coverage classification for a final rendered slide."""
    native_text = [item for item in layout.get("text") or [] if item.get("native")]
    native_images = [item for item in layout.get("images") or [] if item.get("native")]
    unsupported = [
        {"kind": "rendered-text", "selector": item.get("selector"), "text": item.get("text", "")[:120], "reason": item.get("reason") or "not native", "preserved_in": "background"}
        for item in layout.get("text") or [] if not item.get("native")
    ] + [
        {"kind": "rendered-image", "selector": item.get("selector"), "source": item.get("src", ""), "reason": item.get("reason") or "not native", "preserved_in": "background"}
        for item in layout.get("images") or [] if not item.get("native")
    ]
    return {"native_text": native_text, "native_images": native_images, "unsupported": unsupported}


def _px_to_emu(value: float) -> int:
    return round(float(value) * SLIDE_WIDTH_EMU / SLIDE_WIDTH_PX)


def _rgb(value: str) -> tuple[int, int, int] | None:
    values = re.findall(r"[\d.]+", value or "")
    return tuple(max(0, min(255, round(float(item)))) for item in values[:3]) if len(values) >= 3 else None  # type: ignore[return-value]


def _font_name(value: str) -> str:
    return (value or "Arial").split(",", 1)[0].strip().strip("\"'") or "Arial"


def _position_fraction(token: str) -> float:
    token = token.strip().lower()
    if token in {"left", "top"}: return 0.0
    if token in {"right", "bottom"}: return 1.0
    if token.endswith("%"):
        try: return max(0, min(1, float(token[:-1]) / 100))
        except ValueError: pass
    return .5


def _object_position(value: str) -> tuple[float, float]:
    parts = value.split() or ["50%"]
    return _position_fraction(parts[0]), _position_fraction(parts[1] if len(parts) > 1 else parts[0])


def _set_east_asian_font(run: Any, typeface: str) -> None:
    from lxml import etree
    from pptx.oxml.ns import qn
    properties = run._r.get_or_add_rPr(); east_asia = properties.find(qn("a:ea"))
    if east_asia is None: east_asia = etree.SubElement(properties, qn("a:ea"))
    east_asia.set("typeface", typeface)


def _add_text(slide: Any, item: dict) -> None:
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
    from pptx.util import Pt
    box = item["box"]
    shape = slide.shapes.add_textbox(_px_to_emu(box["x"]), _px_to_emu(box["y"]), max(1, _px_to_emu(box["width"])), max(1, _px_to_emu(box["height"])))
    frame = shape.text_frame; frame.clear(); frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0; frame.word_wrap = item.get("whiteSpace") not in {"nowrap", "pre"}; frame.auto_size = MSO_AUTO_SIZE.NONE; frame.vertical_anchor = MSO_ANCHOR.TOP
    paragraph = frame.paragraphs[0]; paragraph.text = item["text"]; paragraph.space_before = paragraph.space_after = Pt(0)
    paragraph.alignment = {"center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT, "end": PP_ALIGN.RIGHT, "justify": PP_ALIGN.JUSTIFY}.get(item.get("textAlign"), PP_ALIGN.LEFT)
    run = paragraph.runs[0]; name = _font_name(item.get("fontFamily", "")); run.font.name = name; _set_east_asian_font(run, name)
    run.font.size = Pt(float(str(item.get("fontSize") or "16").removesuffix("px")) * .75)
    weight = str(item.get("fontWeight") or "400"); run.font.bold = weight == "bold" or (weight.isdigit() and int(weight) >= 600); run.font.italic = item.get("fontStyle") in {"italic", "oblique"}
    if color := _rgb(item.get("color", "")): run.font.color.rgb = RGBColor(*color)
    if (line_height := str(item.get("lineHeight") or "")).endswith("px"): paragraph.line_spacing = Pt(float(line_height[:-2]) * .75)


def _source_to_file(source: str, project: Path, directory: Path, ordinal: int) -> tuple[Path, str | None]:
    """Resolve final-DOM data/local image sources for python-pptx."""
    if source.startswith("data:"):
        header, separator, payload = source.partition(",")
        if not separator or ";base64" not in header.lower(): return Path(), "image data URL is not base64 encoded"
        mime = header[5:].split(";", 1)[0].lower(); suffix = {"image/png": ".png", "image/jpeg": ".jpg", "image/gif": ".gif", "image/webp": ".webp"}.get(mime)
        if not suffix: return Path(), f"{mime or 'unknown'} cannot be embedded reliably by python-pptx"
        target = directory / f"rendered-{ordinal:03d}{suffix}"
        try: target.write_bytes(base64.b64decode(payload, validate=True))
        except ValueError: return Path(), "image data URL is corrupt"
        return target, None
    if source.startswith("file:"):
        from urllib.parse import unquote, urlparse
        path = Path(unquote(urlparse(source).path)).resolve()
        if not path.is_relative_to(project) or not path.is_file(): return Path(), "rendered image is outside the project or missing"
        if path.suffix.lower() not in {".png", ".jpg", ".jpeg", ".gif", ".webp"}: return Path(), f"{path.suffix or 'unknown'} cannot be embedded reliably by python-pptx"
        return path, None
    return Path(), "rendered image is not a local or embedded source"


def _add_picture(slide: Any, image_path: Path, item: dict) -> None:
    box = item["box"]; left, top, width, height = _px_to_emu(box["x"]), _px_to_emu(box["y"]), _px_to_emu(box["width"]), _px_to_emu(box["height"])
    source_width, source_height = max(1, int(item.get("naturalWidth") or 1)), max(1, int(item.get("naturalHeight") or 1))
    fit, x_fraction, y_fraction = item.get("objectFit") or "fill", *_object_position(item.get("objectPosition") or "50% 50%")
    if fit == "contain":
        scale = min(width / source_width, height / source_height); picture_width, picture_height = round(source_width * scale), round(source_height * scale)
        slide.shapes.add_picture(str(image_path), left + round((width - picture_width) * x_fraction), top + round((height - picture_height) * y_fraction), picture_width, picture_height)
    else:
        picture = slide.shapes.add_picture(str(image_path), left, top, width, height)
        if fit == "cover":
            source_ratio, box_ratio = source_width / source_height, width / height
            if source_ratio > box_ratio:
                shown = box_ratio / source_ratio; picture.crop_left = (1 - shown) * x_fraction; picture.crop_right = (1 - shown) * (1 - x_fraction)
            elif source_ratio < box_ratio:
                shown = source_ratio / box_ratio; picture.crop_top = (1 - shown) * y_fraction; picture.crop_bottom = (1 - shown) * (1 - y_fraction)


def _validate_pptx(path: Path, expected_slides: int, pptx: Any) -> None:
    reopened = pptx.Presentation(str(path))
    if len(reopened.slides) != expected_slides or reopened.slide_width != SLIDE_WIDTH_EMU or reopened.slide_height != SLIDE_HEIGHT_EMU: raise RuntimeError("PPTX verification found an invalid slide set or size.")


def _install_pair_atomically(pptx_temp: Path, output: Path, report_temp: Path, report: Path) -> None:
    backups: list[tuple[Path, Path]] = []; installed: list[Path] = []
    try:
        for destination in (output, report):
            if destination.exists():
                fd, name = tempfile.mkstemp(prefix=f".{destination.name}.", suffix=".bak", dir=destination.parent); os.close(fd); backup = Path(name); backup.unlink(); os.replace(destination, backup); backups.append((destination, backup))
        for temporary, destination in ((pptx_temp, output), (report_temp, report)): os.replace(temporary, destination); installed.append(destination)
    except Exception:
        for item in installed: item.unlink(missing_ok=True)
        for destination, backup in reversed(backups): os.replace(backup, destination)
        raise
    finally:
        for _, backup in backups: backup.unlink(missing_ok=True)


def export_hybrid_pptx(*, project: Path, deck: dict, html_path: Path, chrome: str, output: Path, report_output: Path, capture: Callable[[str, Path, dict, Path], list[dict]] = collect_render_layers) -> dict:
    """Build a hybrid PPTX using only the deck manifest and rendered final HTML."""
    pptx = require_python_pptx(); project, html_path, output, report_output = project.resolve(), html_path.resolve(), output.resolve(), report_output.resolve()
    if not html_path.is_file(): raise SystemExit(f"Canonical HTML is missing: {html_path}")
    if output.suffix.lower() != ".pptx" or report_output.suffix.lower() != ".json" or output == report_output: raise SystemExit("PPTX output must be .pptx and coverage output must be a distinct .json file.")
    if not output.parent.is_dir() or not report_output.parent.is_dir(): raise SystemExit("PPTX output parent directories must already exist.")
    paths = deck.get("slides") if isinstance(deck, dict) else None
    if not isinstance(paths, list) or not paths: raise SystemExit("PPTX export requires a deck manifest with slide HTML paths.")
    work = Path(tempfile.mkdtemp(prefix=".oil-ppt-pptx-", dir=project)); fd, name = tempfile.mkstemp(prefix=f".{output.name}.", suffix=".tmp", dir=output.parent); os.close(fd); pptx_temp = Path(name); fd, name = tempfile.mkstemp(prefix=f".{report_output.name}.", suffix=".tmp", dir=report_output.parent); os.close(fd); report_temp = Path(name)
    try:
        backgrounds = work / "backgrounds"; backgrounds.mkdir(); layouts = capture(chrome, html_path, deck, backgrounds)
        if len(layouts) != len(paths): raise RuntimeError("Browser capture returned an incomplete slide set.")
        presentation = pptx.Presentation(); presentation.slide_width = SLIDE_WIDTH_EMU; presentation.slide_height = SLIDE_HEIGHT_EMU; presentation.core_properties.title = str(deck.get("title") or "oil-ppt"); presentation.core_properties.created = presentation.core_properties.modified = datetime(2000, 1, 1, tzinfo=timezone.utc)
        while presentation.slides: presentation.part.drop_rel(presentation.slides._sldIdLst[0].rId); del presentation.slides._sldIdLst[0]
        per_slide: list[dict] = []; total_text = total_media = total_unsupported = 0
        for page, layout in enumerate(layouts, start=1):
            slide = presentation.slides.add_slide(presentation.slide_layouts[6]); slide.shapes.add_picture(layout["background"], 0, 0, width=SLIDE_WIDTH_EMU, height=SLIDE_HEIGHT_EMU)
            classified = classify_rendered_slide(layout); unsupported = list(classified["unsupported"])
            for item in classified["native_text"]: _add_text(slide, item)
            native_media = 0
            for ordinal, item in enumerate(classified["native_images"], start=1):
                image_path, reason = _source_to_file(str(item.get("src") or ""), project, work, page * 1000 + ordinal)
                if reason: unsupported.append({"kind": "rendered-image", "selector": item.get("selector"), "source": item.get("src", ""), "reason": reason, "preserved_in": "background"}); continue
                _add_picture(slide, image_path, item); native_media += 1
            unsupported_count = len(unsupported); per_slide.append({"page": page, "slide_id": layout.get("slideId"), "title": layout.get("title"), "native_text_count": len(classified["native_text"]), "native_media_count": native_media, "rasterized_regions": [{"kind": "deterministic-background", "source": str(html_path), "sha256": _sha256(Path(layout["background"])), "contains": "CSS geometry, decoration, and non-native rendered content"}], "unsupported_structure_count": unsupported_count, "unsupported_structures": unsupported})
            total_text += len(classified["native_text"]); total_media += native_media; total_unsupported += unsupported_count
        presentation.save(str(pptx_temp)); _validate_pptx(pptx_temp, len(paths), pptx)
        theme = validate_theme(deck.get("theme")); palette = {"name": theme["palette"], **PALETTES[theme["palette"]]}; profile = theme["typography"]; typeface = _font_name(TYPOGRAPHY[profile]["font-zh"])
        native_total, denominator = total_text + total_media, total_text + total_media + total_unsupported
        report = {"schema_version": COVERAGE_SCHEMA, "ok": True, "project": str(project), "canonical_html": str(html_path), "pptx": str(output), "pptx_sha256": _sha256(pptx_temp), "slide_size": {"width_emu": SLIDE_WIDTH_EMU, "height_emu": SLIDE_HEIGHT_EMU, "aspect_ratio": "16:9"}, "theme": {"palette": palette, "typography": profile, "typeface": typeface}, "summary": {"slide_count": len(per_slide), "native_text_count": total_text, "native_media_count": total_media, "native_object_count": native_total, "rasterized_region_count": len(per_slide), "unsupported_structure_count": total_unsupported, "coverage_denominator": denominator, "editable_coverage_percent": round(100 * native_total / denominator, 2) if denominator else 100.0}, "methodology": "Coverage counts only final-DOM text and images restored as safe native objects. Every other rendered result remains in a deterministic slide background.", "slides": per_slide}
        report_temp.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"); _install_pair_atomically(pptx_temp, output, report_temp, report_output); return report
    finally:
        pptx_temp.unlink(missing_ok=True); report_temp.unlink(missing_ok=True); shutil.rmtree(work, ignore_errors=True)
